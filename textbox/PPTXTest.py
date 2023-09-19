from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# Create a new presentation
prs = Presentation()

# Add a slide
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Set the position and size of the text box
left = Inches(1)
top = Inches(1)
width = Inches(2)
height = Inches(5)

# Add a text box to the slide
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

# Rotate the text box vertically
txBox.rotation = 270

# Add text to the text box
p = tf.add_paragraph()
p.text = "Vertical Text"
p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT  # Adjust alignment as needed

# Save the presentation
prs.save("vertical_text_example.pptx")
